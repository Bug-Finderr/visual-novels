"""Build a polished deck of the v2 refactor.

Layout: 16:9 widescreen, candlelit-amber theme matching the UI.
~17 slides covering the architecture shift, runtime path, save/restart,
chapter continuation, and the Netlify deployment plan.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent
OUT = DOCS / "output" / "storyplex-v2.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)
DIAG = DOCS / "diagrams" / "v2" / "png"

# Candlelit palette — same vibe as the new UI
BG       = RGBColor(0x07, 0x06, 0x0A)
BG_RAISE = RGBColor(0x14, 0x11, 0x0B)
CARD     = RGBColor(0x1C, 0x16, 0x0E)
CARD_2   = RGBColor(0x18, 0x14, 0x0E)
INK      = RGBColor(0xF3, 0xE7, 0xCF)
MUTED    = RGBColor(0xC4, 0xB0, 0x8A)
DIM      = RGBColor(0x8C, 0x7A, 0x5B)
GOLD     = RGBColor(0xD4, 0xA8, 0x57)
GOLD_BR  = RGBColor(0xEB, 0xC4, 0x73)
EMBER    = RGBColor(0xC4, 0x66, 0x3A)
GREEN    = RGBColor(0x91, 0xA8, 0x7A)
DANGER   = RGBColor(0xC2, 0x5D, 0x52)

SLIDE_W, SLIDE_H = Inches(13.33), Inches(7.5)
SERIF = "Cormorant Garamond"
SANS  = "Inter"
MONO  = "JetBrains Mono"


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, width=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill(bg, BG)
    # Warm radial glow accent at top
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(-2), Inches(-3),
                                   Inches(17), Inches(7))
    fill(glow, RGBColor(0x1C, 0x12, 0x07))
    glow.line.fill.background()


def text_box(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    for i, run in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = run["text"]
        f = r.font
        f.name = run.get("font", SANS)
        f.size = Pt(run.get("size", 18))
        f.bold = run.get("bold", False)
        f.italic = run.get("italic", False)
        if run.get("color") is not None:
            f.color.rgb = run["color"]
        if run.get("space_before") is not None:
            p.space_before = Pt(run["space_before"])
        if run.get("space_after") is not None:
            p.space_after = Pt(run["space_after"])
    return tb


def ornament(slide, x, y, w=Inches(2.0)):
    """Thin gold horizontal rule."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.02))
    fill(bar, GOLD)


def header(slide, title, subtitle=None):
    text_box(
        slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.9),
        [{"text": title, "font": SERIF, "size": 38, "color": INK}],
    )
    ornament(slide, Inches(0.75), Inches(1.4))
    if subtitle:
        text_box(
            slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5),
            [{"text": subtitle, "font": SERIF, "size": 17,
              "italic": True, "color": GOLD_BR}],
        )


_PAGE_STATE = {"current": 0}


def page_no(slide, n=None, total=None):
    """Auto-incrementing page label. Pass nothing — the next slide just calls
    `page_no(s)` and the counter ticks. Total is filled in from the global
    TOTAL constant at render time so we never need to hand-renumber."""
    if n is None:
        _PAGE_STATE["current"] += 1
        n = _PAGE_STATE["current"]
    if total is None:
        total = TOTAL
    text_box(
        slide, Inches(12.3), Inches(7.05), Inches(1.0), Inches(0.3),
        [{"text": f"{n} / {total}", "font": MONO, "size": 9, "color": DIM}],
        align=PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(s, CARD)
    stroke(s, GOLD, 0.5)
    return s


def bullet_block(slide, x, y, w, h, bullets, *, size=15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            # (label, value)
            r1 = p.add_run(); r1.text = item[0] + "  "
            r1.font.name = SANS; r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = GOLD_BR
            r2 = p.add_run(); r2.text = item[1]
            r2.font.name = SANS; r2.font.size = Pt(size); r2.font.color.rgb = INK
        else:
            r0 = p.add_run(); r0.text = "•  "
            r0.font.name = SANS; r0.font.size = Pt(size); r0.font.color.rgb = GOLD
            r1 = p.add_run(); r1.text = item
            r1.font.name = SANS; r1.font.size = Pt(size); r1.font.color.rgb = INK


def mono_block(slide, x, y, w, h, text, *, size=12):
    """Monospace code/pre block on a darker card."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(bg, CARD_2)
    stroke(bg, RGBColor(0x57, 0x49, 0x2D), 0.4)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line if line else " "
        r.font.name = MONO; r.font.size = Pt(size); r.font.color.rgb = INK
        p.space_after = Pt(0)


def table_block(slide, x, y, w, h, rows, *, col_widths=None, header_color=GOLD_BR):
    """Lightweight 2-column or 3-column table via textboxes."""
    n_rows = len(rows)
    row_h = h / n_rows
    ncols = max(len(r) for r in rows)
    if col_widths is None:
        col_widths = [w / ncols] * ncols
    for ri, row in enumerate(rows):
        cy = y + row_h * ri
        cx = x
        for ci, cell in enumerate(row):
            cw = col_widths[ci]
            tb = slide.shapes.add_textbox(cx, cy, cw, row_h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(cell)
            r.font.name = SANS; r.font.size = Pt(14)
            if ri == 0:
                r.font.color.rgb = header_color
                r.font.bold = True
            else:
                r.font.color.rgb = INK
            cx += cw
        # Divider under header row
        if ri == 0:
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y + row_h - Inches(0.02),
                w, Inches(0.01),
            )
            fill(div, RGBColor(0x57, 0x49, 0x2D))


def diagram_slide(slide, png_filename, *, top=Inches(1.85), height=Inches(5.35)):
    """Center a PNG diagram horizontally on a slide below the header."""
    path = DIAG / png_filename
    if not path.exists():
        raise SystemExit(f"missing diagram: {path}")
    # cairosvg writes 1920×1080 (16:9). Fit to slide width minus margins.
    w = Inches(12)
    x = (SLIDE_W - w) // 2
    slide.shapes.add_picture(str(path), x, top, width=w, height=height)


# =====================================================================
# Slides
# =====================================================================

prs = make_prs()
TOTAL = 25  # 19 text slides + 6 diagram slides


def add_title():
    s = blank(prs); set_bg(s)
    text_box(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
             [{"text": "STORYPLEX",
               "font": SERIF, "size": 18, "color": GOLD,
               "bold": False}])
    text_box(s, Inches(0.7), Inches(2.6), Inches(12), Inches(2.2),
             [{"text": "Storyplex — v2",
               "font": SERIF, "size": 84, "color": INK}])
    ornament(s, Inches(0.75), Inches(4.7), Inches(2.4))
    text_box(s, Inches(0.7), Inches(4.85), Inches(12), Inches(1.0),
             [{"text": "Pre-compiled stories, in-house voices.",
               "font": SERIF, "size": 28, "italic": True, "color": MUTED}])
    text_box(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.5),
             [{"text": "A snapshot of the recent architecture refactor.",
               "font": SANS, "size": 16, "color": DIM}])
    page_no(s)


def add_headline_features():
    """Banner slide highlighting the two headline v2 features."""
    s = blank(prs); set_bg(s)
    header(s, "Headline features in v2",
           "Two things that change the whole shape of the product")

    # Left card — Pre-compiled story
    card(s, Inches(0.75), Inches(2.2), Inches(5.85), Inches(4.5))
    text_box(s, Inches(1.0), Inches(2.35), Inches(5.4), Inches(0.5),
             [{"text": "NEW",
               "font": SANS, "size": 11, "bold": True, "color": EMBER}])
    text_box(s, Inches(1.0), Inches(2.7), Inches(5.4), Inches(0.8),
             [{"text": "Pre-compiled story",
               "font": SERIF, "size": 32, "color": GOLD_BR}])
    ornament(s, Inches(1.0), Inches(3.55), Inches(1.4))
    bullet_block(s, Inches(1.0), Inches(3.75), Inches(5.4), Inches(2.8),
                 [
                     "World + all 10 beat dialogues + all 5 endings + every TTS line "
                     "baked into the cache BEFORE the player starts.",
                     "Runtime is a deterministic DB lookup — <100 ms per page.",
                     "No spinner between choices. No 'characters are responding' delay.",
                 ], size=13)

    # Right card — In-house TTS
    card(s, Inches(6.85), Inches(2.2), Inches(5.7), Inches(4.5))
    text_box(s, Inches(7.1), Inches(2.35), Inches(5.4), Inches(0.5),
             [{"text": "NEW",
               "font": SANS, "size": 11, "bold": True, "color": EMBER}])
    text_box(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.8),
             [{"text": "In-house TTS — Silk Mulberry",
               "font": SERIF, "size": 28, "color": GOLD_BR}])
    ornament(s, Inches(7.1), Inches(3.55), Inches(1.4))
    bullet_block(s, Inches(7.1), Inches(3.75), Inches(5.2), Inches(2.8),
                 [
                     "Our own expressive English voice engine — built in-house.",
                     "WebSocket streaming: PCM arrives as it's synthesized.",
                     "Stable per-character description + per-line emotion delta — "
                     "consistent identity across every line.",
                 ], size=13)
    page_no(s)


def add_problem():
    s = blank(prs); set_bg(s)
    header(s, "The problem we set out to fix",
           "Why the old runtime felt sluggish")
    bullet_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3),
                 [
                     "Old runtime called Gemini Flash on every page advance.",
                     "2–4 statements per LLM call → 4–6 calls per beat → 5–8 s of wait.",
                     "Choice buttons hid behind a spinner.",
                     "TTS streamed on demand but spun up a fresh WS per line.",
                 ], size=17)
    card(s, Inches(0.75), Inches(5.4), Inches(11.8), Inches(1.3))
    text_box(s, Inches(1.1), Inches(5.65), Inches(11.0), Inches(1.0),
             [{"text": "“User don’t feel the wait and lag.”",
               "font": SERIF, "size": 26, "italic": True, "color": GOLD_BR}])
    page_no(s)


def add_shift():
    s = blank(prs); set_bg(s)
    header(s, "The shift", "What pre-generation covers now")
    table_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4.4),
                rows=[
                    ["", "Before", "After"],
                    ["Story spine + endings", "pre-gen ✓", "pre-gen ✓"],
                    ["Per-beat dialogue", "live LLM per turn", "pre-gen ✓ + cached"],
                    ["Ending dialogue", "live LLM at beat 9", "pre-gen ✓ + cached"],
                    ["Voice audio", "streamed per line", "pre-gen ✓ + cached"],
                    ["Per-page advance", "1–1.5 s", "<100 ms DB lookup"],
                ],
                col_widths=[Inches(4.2), Inches(3.8), Inches(3.8)])
    text_box(s, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.4),
             [{"text": "Free-input is the only thing that still calls Gemini at runtime.",
               "font": SERIF, "size": 16, "italic": True, "color": MUTED}])
    page_no(s)


def add_story_shape():
    s = blank(prs); set_bg(s)
    header(s, "Story shape", "Spine + branching beats")
    bullet_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3),
                 [
                     "10 beats in a fixed spine.",
                     "3 pre-baked choices per beat, each tagged (alignmentTag + magnitude).",
                     "5 candidate endings keyed by tag.",
                     "Choices don’t change beat content; they nudge alignment_state.",
                     "The ending with the highest score fires at beat 9.",
                 ], size=17)
    text_box(s, Inches(0.75), Inches(5.9), Inches(11.8), Inches(1.5),
             [{"text": "Linear in content. Branching in destination.",
               "font": SERIF, "size": 26, "italic": True, "color": GOLD_BR},
              {"text": "",  "size": 6},
              {"text": "Most players see the same dialogue. Their choices route them to one of 5 distinct epilogues.",
               "font": SANS, "size": 14, "color": MUTED}])
    page_no(s)


def add_spine_flow():
    s = blank(prs); set_bg(s)
    header(s, "Spine flow", "How beats connect")
    mono_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3.5),
"""opening ──► beat 0 ──┬─► beat 1 ──┬─► beat 2 ──► ... ──► beat 9 ──► ending
                     │            │
                  choice        choice
              (3 options,    (3 options,
               each tagged)    each tagged)

         alignment_state at beat 9 picks 1 of 5 cached endings.""",
                size=14)
    bullet_block(s, Inches(0.75), Inches(5.9), Inches(11.8), Inches(1.3),
                 [
                     "Every choice routes to the SAME next beat — only the ending changes.",
                     "_ensure_scene_change + _ensure_cast keep the stage clean across beats.",
                 ], size=14)
    page_no(s)


def add_pipeline():
    s = blank(prs); set_bg(s)
    header(s, "Generation pipeline", "Parallel everywhere")
    table_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4.4),
                rows=[
                    ["Phase", "What runs", "⏱"],
                    ["A", "World + spine + endings (Pro 2.5)", "45 s"],
                    ["B", "Character neutrals in parallel", "12 s"],
                    ["C", "9 emotions × N chars + N scenes pooled", "~100 s"],
                    ["D", "Overlays ││ 8 beat dialogues (Flash)", "~50 s"],
                    ["E", "Script + voices + 5 endings", "~15 s"],
                    ["F", "TTS pre-render of every line (8 workers)", "~1.5–2 min"],
                ],
                col_widths=[Inches(1.4), Inches(8.4), Inches(2.0)])
    text_box(s, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.4),
             [{"text": "Total: ~3.5–5 min. After that — zero waits during play.",
               "font": SERIF, "size": 16, "italic": True, "color": GOLD_BR}])
    page_no(s)


def add_caches():
    s = blank(prs); set_bg(s)
    header(s, "Caches we write", "What lands in the DB + disk")
    table_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4.4),
                rows=[
                    ["Table", "Holds"],
                    ["sessions", "spine, endings, alignment, chosen ending, chapter linkage"],
                    ["characters", "personality, voice_caption, gender (new)"],
                    ["beat_expansions", "ALL 10 beats' statements (pre-rendered)"],
                    ["ending_dialogue", "ALL 5 endings' statements (pre-rendered)"],
                    ["script_labels", "flat statements per label"],
                    ["saves", "Ren'Py-style checkpoints"],
                ],
                col_widths=[Inches(3.4), Inches(8.4)])
    text_box(s, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.4),
             [{"text": "Plus on-disk: sprites, backgrounds, overlays, <sha1>.wav audio.",
               "font": SANS, "size": 14, "color": MUTED}])
    page_no(s)


def add_runtime():
    s = blank(prs); set_bg(s)
    header(s, "Runtime — per page click", "What happens in <100 ms")
    mono_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4.4),
"""client ──► /choice {alignmentTag, magnitude}
            │
            ▼
   process_player_action
            │
            ├── apply alignment to alignment_state[ending_id]
            ├── target_beat = beat_index + 1
            ├── LOOKUP beat_expansions  ◄─── cache hit
            ├── _ensure_scene_change(canonical sceneId)
            ├── _ensure_cast(hide non-cast chars)
            └── append beat's pre-baked choices
            ▼
   200 OK in <100 ms""",
                size=14)
    page_no(s)


def add_tts():
    s = blank(prs); set_bg(s)
    header(s, "TTS — WebSocket all the way",
           "Stream from disk if cached, from Mulberry if not")
    bullet_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4),
                 [
                     "Client opens WS to /api/sessions/{sid}/tts/stream.",
                     "Server checks the on-disk WAV cache.",
                     "Cache hit: streams PCM payload in 100 ms chunks.",
                     "Cache miss: opens upstream Mulberry WS, forwards each PCM frame as it arrives, writes WAV at end.",
                     "Client decodes int16 LE @ 24 kHz, schedules each BufferSource on the first chunk.",
                     "AnalyserNode drives amplitude-based lip-sync.",
                 ], size=15)
    page_no(s)


def add_silk_mulberry():
    """Dedicated headline slide for the in-house TTS engine."""
    s = blank(prs); set_bg(s)
    header(s, "Silk Mulberry — our in-house TTS",
           "Built in-house, streamed over WebSocket, voiced per character")

    # Left: what it is
    card(s, Inches(0.75), Inches(2.2), Inches(5.9), Inches(4.5))
    text_box(s, Inches(1.0), Inches(2.4), Inches(5.5), Inches(0.5),
             [{"text": "What it is",
               "font": SERIF, "size": 22, "italic": True, "color": GOLD_BR}])
    ornament(s, Inches(1.0), Inches(2.95), Inches(1.2))
    bullet_block(s, Inches(1.0), Inches(3.15), Inches(5.5), Inches(3.4),
                 [
                     "Custom expressive English voice model — built in-house.",
                     "POST /v1/tts/ws-connect mints a one-shot WS session.",
                     "Server forwards each PCM frame to the browser AS IT ARRIVES.",
                     "First-byte latency: <500 ms after the click.",
                 ], size=13)

    # Right: voice routing
    card(s, Inches(6.75), Inches(2.2), Inches(5.8), Inches(4.5))
    text_box(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(0.5),
             [{"text": "How a voice is built",
               "font": SERIF, "size": 22, "italic": True, "color": GOLD_BR}])
    ornament(s, Inches(7.0), Inches(2.95), Inches(1.2))
    bullet_block(s, Inches(7.0), Inches(3.15), Inches(5.4), Inches(3.4),
                 [
                     ("base", "per-character voice_caption (stable across every line)"),
                     ("+ delta", "per-line emotion/pacing hint from the expression tag"),
                     ("speaker", "speaker_1 preset for females; description-driven for males"),
                     ("pitch", "f0_up_key shifted by age + gender"),
                 ], size=13)
    page_no(s)


def add_voice_default():
    s = blank(prs); set_bg(s)
    header(s, "Voice profiles", "Stable identity + per-line delta")
    card(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(2.0))
    text_box(s, Inches(1.1), Inches(2.45), Inches(11.0), Inches(0.4),
             [{"text": "Per-character default — set ONCE from voice_caption:",
               "font": SANS, "size": 14, "color": GOLD_BR, "bold": True}])
    text_box(s, Inches(1.1), Inches(2.95), Inches(11.0), Inches(1.2),
             [{"text": "“A teenage girl with a clear, mid-range voice that is often sharp and clipped, but can soften with vulnerability.”",
               "font": SERIF, "size": 17, "italic": True, "color": INK}])
    text_box(s, Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.4),
             [{"text": "Per-line delta — only the delivery hint changes:",
               "font": SANS, "size": 14, "color": GOLD_BR, "bold": True}])
    table_block(s, Inches(0.75), Inches(4.95), Inches(11.8), Inches(2.0),
                rows=[
                    ["expression", "appended"],
                    ["happy",  "Speak cheerfully and a bit quickly, with a smile in the voice."],
                    ["sad",    "Speak softly and slowly, with a downcast tone."],
                    ["angry",  "Speak sharply and with bite, faster paced."],
                ],
                col_widths=[Inches(2.5), Inches(9.3)])
    page_no(s)


def add_voice_routing():
    s = blank(prs); set_bg(s)
    header(s, "Voice routing", "Explicit gender drives the speaker")
    bullet_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3),
                 [
                     ("gender = female", "speaker_1 preset + age-based pitch"),
                     ("gender = male / neutral", "description-driven, no preset, deeper pitch"),
                     ("narrator", "description-driven, f0_up_key = -3"),
                 ], size=17)
    text_box(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(1.5),
             [{"text": "gender is now an EXPLICIT field on every character (set by the world prompt).",
               "font": SANS, "size": 14, "color": INK},
              {"text": "Legacy sessions fall back to keyword detection over voice_caption.",
               "font": SANS, "size": 14, "color": MUTED}])
    page_no(s)


def add_guards():
    s = blank(prs); set_bg(s)
    header(s, "Stage guards", "Two server-side passes per cache hit")
    card(s, Inches(0.75), Inches(2.2), Inches(5.7), Inches(4.5))
    text_box(s, Inches(1.0), Inches(2.4), Inches(5.3), Inches(0.5),
             [{"text": "_ensure_scene_change",
               "font": SERIF, "size": 20, "color": GOLD_BR, "italic": True}])
    bullet_block(s, Inches(1.0), Inches(3.0), Inches(5.3), Inches(3.5),
                 [
                     "Replace lead scene_change with the spine's canonical sceneId",
                     "or prepend one if missing",
                     "fixes 'background went black after a choice'",
                     "(LLM hallucinated a sceneId that didn't exist on disk)",
                 ], size=13)
    card(s, Inches(6.65), Inches(2.2), Inches(5.9), Inches(4.5))
    text_box(s, Inches(6.9), Inches(2.4), Inches(5.5), Inches(0.5),
             [{"text": "_ensure_cast",
               "font": SERIF, "size": 20, "color": GOLD_BR, "italic": True}])
    bullet_block(s, Inches(6.9), Inches(3.0), Inches(5.5), Inches(3.5),
                 [
                     "Prepend hide_character for every char NOT in beat.castIds",
                     "Inserted after the scene_change",
                     "stops prior-beat sprites from lingering",
                     "safe to over-hide — client no-ops it",
                 ], size=13)
    page_no(s)


def add_save_restart():
    s = blank(prs); set_bg(s)
    header(s, "Save / Load / Restart", "Ren'Py-style checkpoints")
    text_box(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.4),
             [{"text": "Triggered by Esc key or the ☰ Menu button.",
               "font": SERIF, "size": 16, "italic": True, "color": MUTED}])
    table_block(s, Inches(0.75), Inches(2.6), Inches(11.8), Inches(4.0),
                rows=[
                    ["Action", "What it does"],
                    ["Save", "Snapshot {label, statementIndex, sceneId, beat, alignment, ending, visibleChars} to saves table"],
                    ["Load", "Rebuild the scene, mount each char, jump label + index"],
                    ["Restart", "Reset alignment + beat. Drop runtime labels. KEEP beat & ending caches so replay is instant."],
                ],
                col_widths=[Inches(2.0), Inches(9.8)])
    page_no(s)


def add_continue():
    s = blank(prs); set_bg(s)
    header(s, "Continue to next chapter",
           "Spawn a child session linked to the parent")
    mono_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3.6),
"""Chapter 1  ────► chosen_ending_id = "noble_sacrifice"
                                  │
                                  ▼
Chapter 2  (new session)
    parent_session_id ─►  Chapter 1
    chapter_number      = 2
    world + cast        = inherited
    spine + endings     = NEW (continues from Ch1's ending)""",
                size=13)
    text_box(s, Inches(0.75), Inches(6.0), Inches(11.8), Inches(1.2),
             [{"text": "build_continuation_prompt forces re-use of character ids and picks up after the parent's chosen ending.",
               "font": SANS, "size": 14, "color": INK}])
    page_no(s)


def add_ui():
    s = blank(prs); set_bg(s)
    header(s, "UI revamp", "Cinematic VN — candlelit amber")
    bullet_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(4.5),
                 [
                     "Candlelit amber palette — gold accent on near-black.",
                     "Cormorant Garamond display serif + Inter body.",
                     "Film grain + soft vignette overlay.",
                     "Glass-card aesthetic with gold rule dividers.",
                     "Per-view refresh: landing, setup wizard, sessions, loading, pause menu.",
                     "Dialogue box gets a soft gold gradient top edge; choice buttons get a gold left bar on hover.",
                 ], size=15)
    page_no(s)


def add_deploy():
    s = blank(prs); set_bg(s)
    header(s, "Deployment plan", "Split hosting: Netlify + Render")
    table_block(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3),
                rows=[
                    ["Layer", "Host", "Why"],
                    ["Frontend (SPA)", "Netlify", "static, free, perfect for Vite output"],
                    ["Backend (FastAPI + WS)", "Render / Fly / Railway", "persistent disk, WS support, long-running"],
                ],
                col_widths=[Inches(3.8), Inches(3.5), Inches(4.5)])
    bullet_block(s, Inches(0.75), Inches(5.4), Inches(11.8), Inches(2.0),
                 [
                     "Netlify Functions can't host the backend: no persistent WS, no persistent disk, hard execution caps.",
                     "Render Starter ($7/mo) with persistent disk handles SQLite + the data/generated/ tree.",
                     "Single-user / dev usage: under $10/mo for hosting. Gemini calls dominate variable cost.",
                 ], size=14)
    page_no(s)


def add_closing():
    s = blank(prs); set_bg(s)
    set_bg(s)
    text_box(s, Inches(0.7), Inches(2.5), Inches(12), Inches(0.5),
             [{"text": "NEXT",
               "font": SERIF, "size": 18, "color": GOLD}])
    text_box(s, Inches(0.7), Inches(3.0), Inches(12), Inches(2.0),
             [{"text": "What’s next",
               "font": SERIF, "size": 72, "color": INK}])
    ornament(s, Inches(0.75), Inches(5.0), Inches(2.4))
    bullet_block(s, Inches(0.75), Inches(5.3), Inches(11.8), Inches(2.0),
                 [
                     "Deploy frontend to Netlify.",
                     "Deploy backend to Render with persistent disk.",
                     "Wire env vars + production CORS.",
                     "Smoke-test fresh-session → play → save → continue.",
                 ], size=17)
    page_no(s)


# ---------------------- diagram slides --------------------------------

def add_diagram_system():
    s = blank(prs); set_bg(s)
    header(s, "System architecture",
           "Browser ↔ FastAPI ↔ Gemini · Mulberry · SQLite · disk")
    diagram_slide(s, "01-system-architecture.png")
    page_no(s)


def add_diagram_spine():
    s = blank(prs); set_bg(s)
    header(s, "Story spine — at a glance",
           "10 beats, 3 choices each, 5 candidate endings")
    diagram_slide(s, "03-spine-flow.png")
    page_no(s)


def add_diagram_pipeline():
    s = blank(prs); set_bg(s)
    header(s, "Generation pipeline — visual",
           "6 phases. Most run in parallel.")
    diagram_slide(s, "02-pipeline-phases.png")
    page_no(s)


def add_diagram_runtime():
    s = blank(prs); set_bg(s)
    header(s, "Runtime cache flow",
           "Choice click → DB lookup → 200 OK in <100 ms")
    diagram_slide(s, "04-runtime-cache-flow.png")
    page_no(s)


def add_diagram_tts():
    s = blank(prs); set_bg(s)
    header(s, "Silk Mulberry — WebSocket flow",
           "Cache-first, live-stream on miss, write WAV on done")
    diagram_slide(s, "05-tts-websocket.png")
    page_no(s)


def add_diagram_chapter():
    s = blank(prs); set_bg(s)
    header(s, "Chapter continuation",
           "Child session inherits world + cast; gets new spine + endings")
    diagram_slide(s, "06-chapter-continuation.png")
    page_no(s)


# Build — order: each text slide is followed by its supporting diagram
# (diagrams appear AFTER the slide that introduces the concept).
for builder in (
    add_title, add_headline_features,
    add_problem, add_diagram_system,
    add_shift,
    add_story_shape, add_diagram_spine,
    add_spine_flow,
    add_pipeline, add_diagram_pipeline,
    add_caches,
    add_runtime, add_diagram_runtime,
    add_tts, add_silk_mulberry, add_diagram_tts,
    add_voice_default, add_voice_routing, add_guards,
    add_save_restart,
    add_continue, add_diagram_chapter,
    add_ui, add_deploy, add_closing,
):
    builder()

prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size // 1024} KB)")
