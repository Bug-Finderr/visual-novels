"""Build a polished, presentation-ready Storyplex deck.

Layout: 16:9 widescreen · dark theme matching the PDF doc.
~16 slides covering: title, exec summary, architecture, story, sprite pivot,
THA3 face pipeline, TTS, runtime loop, frontend anim, GPU infra, perf, DB,
deployment, roadmap, closing.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
DIAGRAMS_PNG = DOCS / "diagrams" / "png"
OUT = DOCS / "output" / "storyplex-design.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Colors
BG       = RGBColor(0x0A, 0x0A, 0x0F)
BG_ALT   = RGBColor(0x14, 0x14, 0x1F)
CARD     = RGBColor(0x1A, 0x1A, 0x2E)
CARD_2   = RGBColor(0x11, 0x11, 0x1E)
TEXT     = RGBColor(0xE8, 0xE8, 0xF0)
MUTED    = RGBColor(0x98, 0x98, 0xB0)
ACCENT   = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_2 = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_3 = RGBColor(0xC4, 0xB5, 0xFD)
GREEN    = RGBColor(0x86, 0xEF, 0xAC)
CYAN     = RGBColor(0x67, 0xE8, 0xF9)

# 16:9 widescreen — 13.33" × 7.5"
SLIDE_W, SLIDE_H = Inches(13.33), Inches(7.5)


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=14, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter",
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT, font="Inter",
                accent_first_word=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        bullet = p.add_run()
        bullet.text = "•  "
        bullet.font.name = font
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = ACCENT_2
        bullet.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def card(slide, x, y, w, h, *, fill=CARD, stroke=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.05
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stroke
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def header(slide, num, title, subtitle=None):
    """Section header strip at the top of a content slide."""
    # Numeric badge
    bx, by = Inches(0.5), Inches(0.35)
    add_text(slide, bx, by, Inches(1.0), Inches(0.45),
             num, size=22, color=ACCENT_2, bold=True)
    add_text(slide, Inches(0.95), Inches(0.32), Inches(11.5), Inches(0.55),
             title, size=24, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.95), Inches(0.78), Inches(11.5), Inches(0.4),
                 subtitle, size=12, color=MUTED)
    # divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(1.20),
                                   Inches(12.33), Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background(); line.shadow.inherit = False


def footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.10), Inches(6), Inches(0.3),
             "Storyplex · Internal · v1.0 · April 2026", size=9, color=MUTED)
    add_text(slide, Inches(7.0), Inches(7.10), Inches(5.83), Inches(0.3),
             f"{page_num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def diagram_slide(prs, page, total, num, title, subtitle, png_path, caption):
    slide = blank(prs)
    set_bg(slide)
    header(slide, num, title, subtitle)

    img_w = Inches(11.5)
    pic = slide.shapes.add_picture(str(png_path), Inches(0.92), Inches(1.55),
                                    width=img_w)
    # Caption under the diagram
    cy = Inches(1.55) + pic.height + Inches(0.10)
    add_text(slide, Inches(0.92), cy, Inches(11.5), Inches(0.4),
             caption, size=10, color=MUTED)
    footer(slide, page, total)
    return slide


# ============================================================
# Build deck
# ============================================================
prs = make_prs()
slides = []

# ---- 1 · Title slide
s = blank(prs); set_bg(s, BG)
# Vertical centering
title_box = s.shapes.add_textbox(Inches(0), Inches(2.1), SLIDE_W, Inches(1.5))
tf = title_box.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Storyplex"
r.font.name = "Inter"; r.font.size = Pt(72); r.font.bold = True
r.font.color.rgb = ACCENT_2
add_text(s, Inches(0), Inches(3.8), SLIDE_W, Inches(0.6),
         "AI-generated visual novels with live characters & Japanese voice acting",
         size=22, color=ACCENT_3, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.0), Inches(4.7), Inches(9.33), Inches(1.2),
         "Story generation · sprite production · GPU character animation (THA3) · "
         "Japanese voice synthesis (Irodori) · runtime gameplay loop",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.6), SLIDE_W, Inches(0.4),
         "INTERNAL  ·  v1.0  ·  APRIL 2026",
         size=10, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
slides.append(s)

# ---- 2 · TOC
s = blank(prs); set_bg(s)
header(s, "", "Contents")
toc_items = [
    ("1", "Executive summary"),
    ("2", "System architecture"),
    ("3", "Tech stack"),
    ("4", "Story generation (Gemini 2.5 Pro)"),
    ("5", "Character sprites — & the pivot"),
    ("6", "Character animation — THA3 face-localized"),
    ("7", "Voice synthesis — Irodori TTS"),
    ("8", "Gameplay loop & runtime AI"),
    ("9", "Frontend animation engine"),
    ("10", "GPU inference infrastructure"),
    ("11", "Performance numbers"),
    ("12", "Data model"),
    ("13", "Deployment & operations"),
    ("14", "Roadmap"),
]
y = Inches(1.6); col_w = Inches(6.0)
for i, (n, t) in enumerate(toc_items):
    col = i // 7
    row = i % 7
    x = Inches(0.7) + col * Inches(6.2)
    yy = y + row * Inches(0.65)
    add_text(s, x, yy, Inches(0.5), Inches(0.5), n, size=18, color=ACCENT_2, bold=True)
    add_text(s, x + Inches(0.55), yy + Inches(0.07), col_w, Inches(0.5), t, size=14, color=TEXT)
slides.append(s)

# ---- 3 · Executive summary
s = blank(prs); set_bg(s)
header(s, "1", "Executive summary",
       "What Storyplex does in one slide")
# Two columns of cards
left_x = Inches(0.5); right_x = Inches(6.92); cw = Inches(5.91); ch = Inches(5.4)
card(s, left_x, Inches(1.5), cw, ch)
add_text(s, left_x + Inches(0.3), Inches(1.65), cw - Inches(0.6), Inches(0.5),
         "What we use", size=15, color=ACCENT_3, bold=True)
add_bullets(s, left_x + Inches(0.3), Inches(2.15), cw - Inches(0.6), Inches(5.0), [
    "Gemini 2.5 Pro — world & opening script (one big JSON, ~30-60s)",
    "Gemini 2.5 Flash Image — 1 neutral sprite per character + 5 backgrounds",
    "THA3 (separable_half, fp16) — face-localized expressions + animations on H100",
    "Irodori-TTS-500M-v2-VoiceDesign — per-line Japanese voice synthesis",
    "Gemini 2.5 Flash — runtime per-turn dialogue",
    "FastAPI + SQLite backend, Vite + vanilla JS frontend, custom VN engine",
], size=13, color=TEXT, line_spacing=1.4)

card(s, right_x, Inches(1.5), cw, ch)
add_text(s, right_x + Inches(0.3), Inches(1.65), cw - Inches(0.6), Inches(0.5),
         "Why it's interesting", size=15, color=ACCENT_3, bold=True)
add_bullets(s, right_x + Inches(0.3), Inches(2.15), cw - Inches(0.6), Inches(5.0), [
    "~10× cheaper image gen — 1 Gemini sprite/char vs the original 10",
    "Face-localized THA3 keeps body Gemini-sharp, only the face animates",
    "Distinct per-character voices via Japanese voice-design captions (gender + age + tone)",
    "Hash-keyed audio manifest — frontend looks up voices deterministically",
    "Runtime TTS arrives in the same response as new dialogue (no second round-trip)",
    "Idempotent, re-runnable pipeline — survives spot preemptions",
], size=13, color=TEXT, line_spacing=1.4)
footer(s, 3, 17)
slides.append(s)

# ---- 4 · System architecture (diagram)
slides.append(diagram_slide(
    prs, 4, 17, "2", "System architecture",
    "Browser ⇄ FastAPI backend ⇄ Gemini API + GPU services on a GCP H100 spot VM",
    DIAGRAMS_PNG / "01-system-architecture.png",
    "Diagram 1 · Three-tier architecture. Backend mediates between the browser and the model providers; "
    "the two GPU services share an H100 in separate uv venvs."))

# ---- 5 · Tech stack table-style
s = blank(prs); set_bg(s)
header(s, "3", "Tech stack")
rows = [
    ("Frontend", "Vite 6 · vanilla ES modules", "Zero framework overhead; custom VN engine"),
    ("Backend", "FastAPI 0.115 · uvicorn · Python 3.13", "Async I/O, Pydantic, hot reload, SSE"),
    ("DB", "SQLite (stdlib)", "Single-process app; file-level backup"),
    ("Story / dialogue", "google-genai SDK · Gemini 2.5 Pro / Flash", "Pro for world, Flash for runtime turns"),
    ("Sprites / BG", "Gemini 2.5 Flash Image", "Native multimodal image gen"),
    ("BG removal", "rembg (U²-Net)", "Reliable transparent cutouts on stylized art"),
    ("Animation", "Talking Head Anime 3 · separable_half · fp16", "12-22 ms / frame on H100; 45-dim pose vector"),
    ("TTS", "Aratako/Irodori-TTS-500M-v2-VoiceDesign · bf16", "RF-DiT 500M + DACVAE; voice-design captions"),
    ("Image proc", "Pillow 12", "Face crop + alpha composite; PNG re-encode"),
    ("HTTP client", "httpx", "Backend → puppeteer / TTS"),
    ("GPU venv", "uv + Python 3.10", "Isolated venvs per service; reproducible installs"),
    ("Cloud", "GCP CE spot VM · 1× H100 80GB · Ubuntu 22.04", "Idle between sessions"),
]
# Header row
hy = Inches(1.5)
add_text(s, Inches(0.5),  hy, Inches(2.7), Inches(0.4), "Layer",       size=12, color=ACCENT_3, bold=True)
add_text(s, Inches(3.3),  hy, Inches(4.5), Inches(0.4), "Tool / version", size=12, color=ACCENT_3, bold=True)
add_text(s, Inches(7.95), hy, Inches(4.9), Inches(0.4), "Why",        size=12, color=ACCENT_3, bold=True)
# divider
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.85), Inches(12.33), Emu(6350))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); ln.shadow.inherit = False
# rows
ry = Inches(1.95)
for layer, tool, why in rows:
    add_text(s, Inches(0.5),  ry, Inches(2.7), Inches(0.42), layer, size=10.5, color=TEXT, bold=True)
    add_text(s, Inches(3.3),  ry, Inches(4.5), Inches(0.42), tool,  size=10.5, color=ACCENT_2)
    add_text(s, Inches(7.95), ry, Inches(4.9), Inches(0.42), why,   size=10.5, color=MUTED)
    ry = ry + Inches(0.40)
footer(s, 5, 17)
slides.append(s)

# ---- 6 · Generation pipeline (diagram)
slides.append(diagram_slide(
    prs, 6, 17, "4", "Generation pipeline",
    "POST /api/sessions/<id>/generate · ~5 minutes total · progress streamed via SSE",
    DIAGRAMS_PNG / "02-generation-pipeline.png",
    "Diagram 2 · Six-step pipeline. Cost-heavy steps (Gemini Image, story Pro) are bracketed by cheap GPU steps."))

# ---- 7 · The pivot
s = blank(prs); set_bg(s)
header(s, "5", "Character sprites — and the pivot",
       "From 10 Gemini Image calls per character to 1; derive the rest with THA3")

# Before / After cards
add_text(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
         "Before", size=15, color=MUTED, bold=True)
card(s, Inches(0.5), Inches(1.95), Inches(6), Inches(2.6))
add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(2.4), [
    "10 Gemini Image calls per character (one per emotion)",
    "Each call ~25-30 s + ~10 s rembg",
    "Drift between expressions: clothing/hair sometimes inconsistent",
    "~$0.5-1 per character · ~5 min for a session",
], size=12, color=TEXT, line_spacing=1.35)

add_text(s, Inches(6.85), Inches(1.5), Inches(6), Inches(0.4),
         "After", size=15, color=GREEN, bold=True)
card(s, Inches(6.85), Inches(1.95), Inches(6), Inches(2.6), stroke=GREEN)
add_bullets(s, Inches(7.05), Inches(2.1), Inches(5.6), Inches(2.4), [
    "1 Gemini Image call per character (neutral only)",
    "9 expressions + 30 animation frames derived via THA3 pose deltas",
    "Guaranteed consistency — every frame is a warp of the same source",
    "~10× cheaper · seconds-per-character on the H100 batch endpoint",
], size=12, color=TEXT, line_spacing=1.35)

# Tradeoff strip
card(s, Inches(0.5), Inches(4.85), Inches(12.33), Inches(1.85), fill=CARD_2, stroke=ACCENT_2)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
         "Tradeoff", size=14, color=ACCENT_3, bold=True)
add_text(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.2),
         "THA3 only does facial-muscle deformations (eyebrows, eye shape, mouth, head pose, iris). "
         "Extreme dramatic moments — tears streaming, body trembling, dramatic hair shifts — fall outside "
         "its capability. For those, we can spawn a one-off Gemini Image call on demand. The current 10 "
         "named emotions cover the vast majority of in-game beats.",
         size=12, color=TEXT, line_spacing=1.4)
footer(s, 7, 17)
slides.append(s)

# ---- 8 · Frame tree (diagram)
slides.append(diagram_slide(
    prs, 8, 17, "5", "Per-character asset tree",
    "1 Gemini sprite + 39 derived frames + N voices",
    DIAGRAMS_PNG / "03-frame-tree.png",
    "Diagram 3 · Each character gets 40 PNGs total: 1 Gemini neutral (green) + 39 THA3-derived face composites (purple)."))

# ---- 9 · THA3 face pipeline (diagram)
slides.append(diagram_slide(
    prs, 9, 17, "6", "THA3 — face-localized animation",
    "Crop face → THA3 inference → composite back over pristine body",
    DIAGRAMS_PNG / "04-tha3-face-pipeline.png",
    "Diagram 4 · Body stays Gemini-sharp; only the face is THA3-warped. Solves THA3's quality drop on full-body inputs."))

# ---- 10 · TTS pipeline (diagram)
slides.append(diagram_slide(
    prs, 10, 17, "7", "Voice synthesis — Irodori TTS",
    "Aratako/Irodori-TTS-500M-v2-VoiceDesign · RF-DiT 500M + DACVAE · Japanese-only",
    DIAGRAMS_PNG / "05-tts-pipeline.png",
    "Diagram 5 · Each TTS call takes JP text + JP voice-design caption. Inference is ~7× real-time on H100."))

# ---- 11 · Runtime loop (diagram)
slides.append(diagram_slide(
    prs, 11, 17, "8", "Gameplay loop & runtime AI",
    "Player choice → Gemini Flash → script_builder → TTS → frontend renders",
    DIAGRAMS_PNG / "07-runtime-loop.png",
    "Diagram 7 · Single round-trip per choice. TTS for new dialogue runs synchronously; manifest delta returned in same response."))

# ---- 12 · Frontend anim stack (diagram)
slides.append(diagram_slide(
    prs, 12, 17, "9", "Frontend animation engine",
    "AnimatedSprite class · 4-frame stack with instant swap · idle breath + blink + mouth-flap",
    DIAGRAMS_PNG / "06-frontend-anim-stack.png",
    "Diagram 6 · Per-character DOM stack. Frame swaps are GPU-composited opacity flips with transition: none."))

# ---- 13 · GPU infra
s = blank(prs); set_bg(s)
header(s, "10", "GPU inference infrastructure",
       "Two services on one H100, separate uv venvs, idempotent bootstrap")
# 3 cards in a row
cw = Inches(4.0); ch = Inches(5.4); cy = Inches(1.5)
gx = [Inches(0.5), Inches(4.67), Inches(8.83)]

card(s, gx[0], cy, cw, ch)
add_text(s, gx[0]+Inches(0.25), cy+Inches(0.18), cw-Inches(0.5), Inches(0.4),
         "VM", size=14, color=ACCENT_3, bold=True)
add_bullets(s, gx[0]+Inches(0.25), cy+Inches(0.6), cw-Inches(0.5), ch-Inches(0.8), [
    "GCP Compute Engine spot VM",
    "1× NVIDIA H100 80GB",
    "driver 580 / CUDA 13",
    "Ubuntu 22.04",
    "2× local NVMe SSD (375 GB each, ephemeral)",
    "Persistent boot disk on the third NVMe",
    "External IP exposed for puppeteer + TTS",
], size=12, color=TEXT, line_spacing=1.4)

card(s, gx[1], cy, cw, ch)
add_text(s, gx[1]+Inches(0.25), cy+Inches(0.18), cw-Inches(0.5), Inches(0.4),
         "Bootstrap (idempotent)", size=14, color=ACCENT_3, bold=True)
add_bullets(s, gx[1]+Inches(0.25), cy+Inches(0.6), cw-Inches(0.5), ch-Inches(0.8), [
    "Mount unformatted NVMe at /workspace",
    "Install uv to ~/.local/bin",
    "Create per-service Python 3.10 venv",
    "Clone source repos (THA3, Irodori-TTS)",
    "Install per-service deps (PyTorch+CUDA, etc.)",
    "Download weights (~860 MB + ~2 GB)",
    "Generate run_*.sh runners",
    "Re-runnable on every spot preemption",
], size=12, color=TEXT, line_spacing=1.4)

card(s, gx[2], cy, cw, ch)
add_text(s, gx[2]+Inches(0.25), cy+Inches(0.18), cw-Inches(0.5), Inches(0.4),
         "Networking", size=14, color=ACCENT_3, bold=True)
add_bullets(s, gx[2]+Inches(0.25), cy+Inches(0.6), cw-Inches(0.5), ch-Inches(0.8), [
    "VM tagged http-server",
    "Existing rule allow-tts-8000 covers puppeteer port",
    "New rule allow-puppeteer-tts-8001 added for TTS",
    "Both bind 0.0.0.0; reachable on external IP",
    "Backend reads PUPPETEER_URL / TTS_URL from .env",
    "httpx client with 60-120s timeouts",
    "(For prod: restrict source ranges or front w/ auth proxy)",
], size=12, color=TEXT, line_spacing=1.4)
footer(s, 13, 17)
slides.append(s)

# ---- 14 · Performance numbers
s = blank(prs); set_bg(s)
header(s, "11", "Performance numbers")
rows = [
    ("Story / opening script", "Gemini 2.5 Pro",                    "~30-60 s",        "Single ~30 KB JSON; response_mime_type=json"),
    ("Per-character sprite",   "Gemini 2.5 Flash Image",            "~25 s + ~10 s",   "rembg post-process; one call per character"),
    ("Per-emotion frame batch","THA3 (separable_half, fp16)",       "~250 ms / 4-frame batch", "12-22 ms/frame raw; one batch per emotion"),
    ("Per-scene background",   "Gemini 2.5 Flash Image",            "~10-15 s",        "5 backgrounds per session"),
    ("Per-line voice",         "Irodori (bf16, 40 steps)",          "~1 s for ~7 s audio", "~7× real-time on H100; 48 kHz mono PCM"),
    ("Per-turn runtime",       "Gemini Flash + Irodori",            "~3-5 s",          "Includes TTS for ~5-15 new lines"),
    ("End-to-end session",     "All combined",                      "~5 min",          "4-character session, generation to game-ready"),
]
hy = Inches(1.5)
add_text(s, Inches(0.5), hy, Inches(3.0), Inches(0.4), "Step",      size=12, color=ACCENT_3, bold=True)
add_text(s, Inches(3.55), hy, Inches(3.7), Inches(0.4), "Model",     size=12, color=ACCENT_3, bold=True)
add_text(s, Inches(7.40), hy, Inches(2.4), Inches(0.4), "Time",      size=12, color=ACCENT_3, bold=True)
add_text(s, Inches(9.95), hy, Inches(3.0), Inches(0.4), "Notes",     size=12, color=ACCENT_3, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.85), Inches(12.33), Emu(6350))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); ln.shadow.inherit = False
ry = Inches(2.0)
for step, mdl, t, note in rows:
    add_text(s, Inches(0.5), ry, Inches(3.0), Inches(0.5), step, size=11, color=TEXT, bold=True)
    add_text(s, Inches(3.55), ry, Inches(3.7), Inches(0.5), mdl,  size=11, color=ACCENT_2)
    add_text(s, Inches(7.40), ry, Inches(2.4), Inches(0.5), t,    size=11, color=GREEN, bold=True)
    add_text(s, Inches(9.95), ry, Inches(3.0), Inches(0.7), note, size=10, color=MUTED)
    ry = ry + Inches(0.55)
# warmup note
card(s, Inches(0.5), Inches(6.05), Inches(12.33), Inches(0.95), fill=CARD_2, stroke=GREEN)
add_text(s, Inches(0.7), Inches(6.18), Inches(12), Inches(0.4),
         "Cold vs warm", size=12, color=GREEN, bold=True)
add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
         "Both GPU services pre-warm at startup. First request is ~770 ms (THA3) / ~1.2 s (Irodori) due to "
         "CUDA kernel JIT — the steady-state numbers above kick in from the second request.",
         size=11, color=TEXT)
footer(s, 14, 17)
slides.append(s)

# ---- 15 · DB schema (diagram)
slides.append(diagram_slide(
    prs, 15, 17, "12", "Data model",
    "5 tables in SQLite + a sibling filesystem tree under data/generated/<sid>/",
    DIAGRAMS_PNG / "08-db-schema.png",
    "Diagram 8 · ER diagram. All child tables CASCADE on session delete. Schema migrations are defensive ALTER TABLEs at startup."))

# ---- 16 · Deployment + Roadmap (combined)
s = blank(prs); set_bg(s)
header(s, "13/14", "Deployment & roadmap")
# Left: deployment
card(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(5.4))
add_text(s, Inches(0.7), Inches(1.65), Inches(6), Inches(0.4),
         "Local dev (one machine)", size=14, color=ACCENT_3, bold=True)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(2.0), [
    "pip install -r server/requirements.txt",
    "Set GEMINI_API_KEY, PUPPETEER_URL, TTS_URL in repo .env",
    "npm run dev:server  →  uvicorn :3001",
    "npm run dev:client  →  vite :3000",
    "Open http://localhost:3000",
], size=12, color=TEXT, line_spacing=1.4)

add_text(s, Inches(0.7), Inches(4.35), Inches(6), Inches(0.4),
         "VM bootstrap (per-incarnation)", size=14, color=ACCENT_3, bold=True)
add_bullets(s, Inches(0.7), Inches(4.8), Inches(5.7), Inches(2.0), [
    "scp + run bootstrap_inference_vm.sh",
    "scp + run bootstrap_tts_vm.sh",
    "nohup /workspace/run_inference.sh & disown",
    "nohup /workspace/run_tts.sh & disown",
    "Curl /health on both ports to verify",
], size=12, color=TEXT, line_spacing=1.4)

# Right: roadmap
card(s, Inches(6.73), Inches(1.5), Inches(6.1), Inches(5.4), stroke=GREEN)
add_text(s, Inches(6.93), Inches(1.65), Inches(6), Inches(0.4),
         "Roadmap", size=14, color=GREEN, bold=True)
add_bullets(s, Inches(6.93), Inches(2.1), Inches(5.7), Inches(4.7), [
    "Anime-trained face detector → better face-crop accuracy",
    "Lip-sync mouth shape to actual TTS audio amplitude",
    "Per-language localization · multilingual TTS",
    "Speculative TTS for likely next dialogue lines (idle prefetch)",
    "GCS-backed asset store for cross-VM session continuity",
    "IAP tunnel auth instead of public 0.0.0.0/0",
    "Save / load slots; export run as a scene-reel video",
    "Player-defined character voices via reference audio",
], size=12, color=TEXT, line_spacing=1.4)

footer(s, 16, 17)
slides.append(s)

# ---- Iteration delta · before vs after this round
s = blank(prs); set_bg(s)
header(s, "★", "What changed in this iteration",
       "Last iteration shipped a Gemini-driven VN on Monogatari + Node. This round we rewrote the stack and added live characters + voice.")

prev_items = [
    "Frontend: Monogatari VN engine (heavy framework, opinionated)",
    "Backend: Node.js + Express + better-sqlite3 / sql.js",
    "Story: Gemini Pro (one-shot world + opening script)",
    "Sprites: 10 Gemini Image calls per character (one per emotion)",
    "Backgrounds: Gemini Image",
    "Animation: none — characters were static PNGs",
    "Voice: none — silent typewriter only",
    "Inference infra: none (no GPU pieces)",
]
new_items = [
    "Frontend: vanilla JS + custom VN engine (~15 files; full control)",
    "Backend: FastAPI + uvicorn + stdlib sqlite3 (httpx for GPU services)",
    "Story: Gemini Pro — now also emits per-line jp text + JP voiceCaption",
    "Sprites: 1 Gemini neutral / character — 9 expressions derived via THA3",
    "Backgrounds: Gemini Image (unchanged) + ken-burns motion in browser",
    "Animation: THA3 puppeteer on H100 — face-localized composite over body",
    "Voice: Irodori-TTS (RF-DiT 500M) on H100 — per-character JP voices",
    "Inference infra: GCP H100 spot VM, two uv venvs, idempotent bootstrap",
]

# Two columns
add_text(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(0.45),
         "Last iteration", size=15, color=MUTED, bold=True)
card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.95), stroke=MUTED)
add_bullets(s, Inches(0.7), Inches(2.15), Inches(5.7), Inches(4.7),
            prev_items, size=11.5, color=MUTED, line_spacing=1.4)

add_text(s, Inches(6.73), Inches(1.5), Inches(6.1), Inches(0.45),
         "This iteration", size=15, color=GREEN, bold=True)
card(s, Inches(6.73), Inches(2.0), Inches(6.1), Inches(4.95), stroke=GREEN)
add_bullets(s, Inches(6.93), Inches(2.15), Inches(5.7), Inches(4.7),
            new_items, size=11.5, color=TEXT, line_spacing=1.4)

footer(s, 17, 18)
slides.append(s)

# ---- 18 · Closing
s = blank(prs); set_bg(s, BG)
add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.5),
         "Thanks", size=64, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.6),
         "Storyplex · technical design · v1.0",
         size=18, color=ACCENT_3, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.0), Inches(9.33), Inches(1.0),
         "8 hand-drawn diagrams · 14 sections · 18 slides · all sources in docs/",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.6), SLIDE_W, Inches(0.4),
         "Questions?  ·  Open the deck PDF for the full write-up",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
slides.append(s)

# Save
prs.save(str(OUT))
print(f"saved: {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")
