"""Approximate-render picture slides from their real coordinates.

Exists because there's no PowerPoint or LibreOffice on this machine — this
catches overlaps, overflow and wrong aspect ratios that a structural check
cannot see. Not pixel-accurate; it does not render theme backgrounds or fonts.
"""
import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation

DOCS = Path(__file__).resolve().parent
SC = 130  # px per inch


def render(deck=DOCS / "output" / "StoryPlex_Capstone.pptx", slides=(8, 9, 12)):
    prs = Presentation(str(deck))
    W, H = int(prs.slide_width / 914400 * SC), int(prs.slide_height / 914400 * SC)
    out = []
    for idx in slides:
        s = prs.slides[idx - 1]
        canvas = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(canvas)
        for sh in s.shapes:
            x, y = int(sh.left / 914400 * SC), int(sh.top / 914400 * SC)
            w, h = max(int(sh.width / 914400 * SC), 1), max(int(sh.height / 914400 * SC), 1)
            if sh.shape_type == 13:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB").resize((w, h))
                canvas.paste(im, (x, y))
                d.rectangle([x, y, x + w, y + h], outline="#999", width=1)
            elif sh.has_text_frame and sh.text_frame.text.strip():
                d.rectangle([x, y, x + w, y + h], outline="#d8d8e0", width=1)
                d.text((x + 4, y + 3), sh.text_frame.text.strip()[:70], fill="#333")
        p = DOCS / "output" / f"preview-slide-{idx}.png"
        canvas.save(p)
        out.append(p)
    return out


if __name__ == "__main__":
    want = [int(a) for a in sys.argv[1:]] or [8, 9, 12]
    for p in render(slides=want):
        print(f"  {p.name}")
