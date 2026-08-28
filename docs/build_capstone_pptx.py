"""Rewrite the BITS capstone template with StoryPlex content.

Edits the supplied deck IN PLACE rather than building a new one, so the
institution's layouts, theme, fonts and master slides are preserved exactly.
Only text content and images change.

Screenshots are picked up from docs/screenshots/ by filename (see the README
there). Anything missing gets a labelled placeholder, so the deck can be built
now and rebuilt as screenshots arrive.

Run:  cd docs && ../server/.venv/bin/python build_capstone_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
TEMPLATE = DOCS / "final_capstone (1).pptx"
SHOTS = DOCS / "screenshots"
CHARTS = DOCS / "charts"
CACHE = DOCS / "output" / ".img-cache"
OUT = DOCS / "output" / "StoryPlex_Capstone.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x6B, 0x6B, 0x80)
PLACEHOLDER_BG = RGBColor(0xEE, 0xF0, 0xF6)


# ---------------------------------------------------------------- helpers
def set_title(slide, text):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0:
            sh.text_frame.text = text
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
            return


def set_body(slide, blocks, size=13):
    """blocks = [(text, level, bold), ...]. Reuses the template's body
    placeholder so theme fonts and colours are inherited."""
    body = None
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 1:
            body = sh
            break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (text, level, bold) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for r in p.runs:
            r.font.size = Pt(size + (1 if bold and level == 0 else 0))
            r.font.bold = bold
            r.font.color.rgb = INK if bold else MUTED


def clear_pictures(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # PICTURE
            sh._element.getparent().remove(sh._element)


def clear_all(slide):
    for sh in list(slide.shapes):
        if not (sh.is_placeholder and sh.placeholder_format.idx in (0, 1)):
            sh._element.getparent().remove(sh._element)


def fit(img_w, img_h, box_w, box_h):
    """Scale to fit inside the box, preserving aspect ratio."""
    scale = min(box_w / img_w, box_h / img_h)
    return img_w * scale, img_h * scale


def place(slide, path, left, top, box_w, box_h, label=None):
    """Drop an image centred in a box. If it's missing, draw a labelled
    placeholder instead so the slot is visible rather than silently empty."""
    from PIL import Image
    if path and Path(path).exists():
        with Image.open(path) as im:
            w, h = im.size
        dw, dh = fit(w, h, box_w, box_h)

        # Retina captures are ~3400px wide; a 10in slide never needs more than
        # ~200 DPI. Downscale into a cache so the deck stays a sane size —
        # 26 MB of source PNGs becomes a few MB with no visible loss.
        target_px = int(dw / 914400 * 200)
        src = Path(path)
        if w > target_px * 1.15:
            CACHE.mkdir(exist_ok=True)
            cached = CACHE / f"{src.stem}@{target_px}.png"
            if not cached.exists():
                with Image.open(src) as im:
                    im.convert("RGB").resize(
                        (target_px, int(h * target_px / w)), Image.LANCZOS
                    ).save(cached, "PNG", optimize=True)
            src = cached

        px, py = int(left + (box_w - dw) / 2), int(top + (box_h - dh) / 2)
        slide.shapes.add_picture(str(src), Emu(px), Emu(py), Emu(int(dw)), Emu(int(dh)))
        return (px, py, int(dw), int(dh))

    # Placeholder matches the 16:9 shape of a real capture, so a missing shot
    # doesn't distort the layout it will eventually occupy.
    ph_h = min(box_h, box_w * 9 / 16)
    ph_y = int(top + (box_h - ph_h) / 2)
    box = slide.shapes.add_textbox(Emu(int(left)), Emu(ph_y), Emu(int(box_w)), Emu(int(ph_h)))
    box.fill.solid(); box.fill.fore_color.rgb = PLACEHOLDER_BG
    box.line.color.rgb = MUTED
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = f"[ screenshot ]\n{label or Path(str(path)).name}"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(12); r.font.color.rgb = MUTED; r.font.bold = True
    return (int(left), ph_y, int(box_w), int(ph_h))


def caption(slide, text, left, top, width, size=11):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Inches(0.3))
    tf = tb.text_frame; tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK
    return tb


# Accept several spellings per slot, so screenshots don't have to be renamed to
# match. First existing match wins; order is preference order.
ALIASES = {
    "01-landing.png":     ["01-landing.png", "landing.png", "home.png"],
    "02-create-form.png": ["02-create-form.png", "create_form.png", "create-form.png", "create.png"],
    "03-loading.png":     ["03-loading.png", "loading.png", "generating.png", "progress.png"],
    "04-credits.png":     ["04-credits.png", "credits.png", "billing.png"],
    "05-reader.png":      ["05-reader.png", "reader.png", "game.png", "play.png"],
    "06-choices.png":     ["06-choices.png", "choices.png", "choice.png"],
    "07-explore.png":     ["07-explore.png", "explore.png", "feed.png"],
    # No library shot supplied; the story detail page fills the slot well.
    "08-library.png":     ["08-library.png", "library.png", "story_page.png", "story-page.png"],
}


def shot(name):
    for candidate in ALIASES.get(name, [name]):
        p = SHOTS / candidate
        if p.exists():
            return p
    return SHOTS / name


# ---------------------------------------------------------------- content
def build():
    prs = Presentation(str(TEMPLATE))
    S = prs.slides
    W, H = prs.slide_width, prs.slide_height
    missing = []

    # ---- 1 title -------------------------------------------------------
    for sh in S[0].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Hawkslab":
            sh.text_frame.text = "StoryPlex"
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
    tb = S[0].shapes.add_textbox(Inches(0.65), Inches(3.35), Inches(5.2), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    tf.text = "AI-generated visual novels"
    p2 = tf.add_paragraph(); p2.text = "storyplex.app"
    for i, p in enumerate(tf.paragraphs):
        for r in p.runs:
            r.font.size = Pt(16 if i == 0 else 12)
            r.font.color.rgb = INK if i == 0 else MUTED
            r.font.bold = i == 0

    # ---- 2 problem statement -------------------------------------------
    set_title(S[1], "Problem Statement")
    set_body(S[1], [
        ("Background of the problem", 0, True),
        ("Visual novels are among the most approachable forms of interactive fiction, but "
         "producing one demands writing, character art, background art and voice acting. A single "
         "short title is months of work across several disciplines, which puts the form out of "
         "reach of almost everyone who would like to tell a story in it.", 0, False),
        ("Gap in existing systems", 0, True),
        ("Generative AI can now produce prose, illustration and speech individually, but the tools "
         "remain separate. Assembling them into a playable, internally consistent story — one where "
         "the same character looks and sounds the same across every scene, and choices actually "
         "change the outcome — is still manual integration work.", 0, False),
        ("Importance of the problem", 0, True),
        ("A system that takes a premise and returns a complete, playable visual novel collapses "
         "months of multi-disciplinary production into minutes, and makes the medium available to "
         "anyone who can describe a story.", 0, False),
    ], size=12)

    # ---- 3 objectives & scope ------------------------------------------
    set_title(S[2], "Objectives & Scope")
    set_body(S[2], [
        ("Objectives:", 0, True),
        ("Objective 1 — Generate a complete, internally consistent visual novel from a short user "
         "premise: plot, cast, branching script, artwork and voices.", 0, False),
        ("Objective 2 — Maintain narrative coherence across a 10-beat story spine with five "
         "distinct endings, where player choices carry weight toward the ending reached.", 0, False),
        ("Objective 3 — Keep generation economically viable by producing only the assets a story "
         "actually uses, rather than a fixed catalogue.", 0, False),
        ("Objective 4 — Operate reliably as a deployed, multi-user web service with authentication, "
         "prepaid billing, and admission control under constrained memory.", 0, False),
        ("Scope:", 0, True),
        ("In scope — story/dialogue generation, character sprites and backgrounds, text-to-speech, "
         "branching gameplay, accounts, publishing and social features, credit-based billing, "
         "production deployment.", 0, False),
        ("Out of scope — video or animation beyond a 2D sprite rig, multiplayer, user-uploaded "
         "artwork, mobile applications, and translation into other languages.", 0, False),
    ], size=11.5)

    # ---- 4 literature review -------------------------------------------
    set_title(S[3], "Existing System / Literature Review")
    set_body(S[3], [
        ("Existing approach 1:", 0, True),
        ("Traditional visual novel engines (Ren'Py, TyranoBuilder) — mature runtimes for playing a "
         "visual novel, but they author nothing. Every line, sprite and background must be supplied "
         "by the creator.", 0, False),
        ("Existing approach 2:", 0, True),
        ("AI story generators (AI Dungeon, NovelAI) — generate prose interactively, but produce no "
         "consistent cast, no artwork tied to the narrative, and no fixed structure, so sessions "
         "drift and rarely reach a coherent ending.", 0, False),
        ("Existing approach 3:", 0, True),
        ("Standalone generative tools (image models, TTS services) used manually and assembled by "
         "hand in an engine.", 0, False),
        ("Limitations:", 0, True),
        ("1) Engines require the creator to already possess writing, art and voice skills.", 0, False),
        ("2) Story generators lose character and world consistency over a long session, and offer "
         "no authored structure or endings.", 0, False),
        ("3) Manual assembly gives no guarantee a character looks the same twice, and the cost of "
         "generating every possible asset is prohibitive.", 0, False),
    ], size=11)

    # ---- 5 architecture -------------------------------------------------
    set_title(S[4], "Proposed System Architecture")
    clear_pictures(S[4])
    set_body(S[4], [
        ("System overview: A React single-page application talks to a FastAPI backend on Render. "
         "The backend orchestrates Gemini for text and images and Silk/Mulberry for speech, "
         "persists to PostgreSQL, and writes generated assets to Google Cloud Storage — from where "
         "the browser loads them directly, never through the server.", 0, False),
    ], size=11)
    place(S[4], CHARTS / "architecture.png",
          Inches(0.45), Inches(2.55), W - Inches(0.9), Inches(4.5))

    # ---- 6 tools & technologies -----------------------------------------
    set_title(S[5], "Tools & Technologies")
    clear_pictures(S[5])
    set_body(S[5], [
        ("Programming Languages:", 0, True),
        ("Python 3.14 (backend) · JavaScript / JSX (frontend)", 0, False),
        ("Frameworks:", 0, True),
        ("1) FastAPI — async web framework, 40 REST endpoints", 0, False),
        ("2) React 18 + Vite — single-page application and build tooling", 0, False),
        ("3) LangGraph — multi-agent story generation pipeline", 0, False),
        ("4) SQLAlchemy + Alembic — data access and schema migrations", 0, False),
        ("Databases & Storage:", 0, True),
        ("1) PostgreSQL — 22 tables, 5 migrations", 0, False),
        ("2) Google Cloud Storage — generated sprites, backgrounds and audio", 0, False),
        ("AI Services:", 0, True),
        ("1) Gemini 3 Flash — story bible generation", 0, False),
        ("2) Gemini 3.1 Flash-Lite — dialogue and beat expansion", 0, False),
        ("3) Gemini 3.1 Flash-Lite Image — sprites, backgrounds, cover art", 0, False),
        ("4) Silk / Mulberry — streaming voice synthesis", 0, False),
        ("Infrastructure:", 0, True),
        ("Render (Docker web service, static site, managed Postgres) · Cashfree Payments · "
         "Google OAuth 2.0 · Docker Compose · Git / GitHub", 0, False),
    ], size=10.5)

    # ---- 7 implementation -----------------------------------------------
    set_title(S[6], "Implementation / Demo")
    set_body(S[6], [
        ("Feature 1 — Multi-agent story generation", 0, True),
        ("A LangGraph pipeline of plot, world, character and chapter agents, gated by a Memory "
         "critic that checks structural consistency and sends work back for revision. Produces a "
         "10-beat spine with five endings and an alignment-weighted choice model.", 0, False),
        ("Feature 2 — Selective asset generation", 0, True),
        ("Dialogue text is generated before any image. The resulting script is then scanned for the "
         "character expressions and scenes it actually references, and only those are drawn — "
         "51% fewer images, with nothing lost that a player would ever see.", 0, False),
        ("Feature 3 — Chroma-key sprite cutout", 0, True),
        ("The image model is asked for a flat magenta background, removed afterwards by adaptive "
         "border-median colour-distance thresholding in numpy. Replacing an ML background remover "
         "with arithmetic cut the memory baseline from 750 MB to 129 MB.", 0, False),
        ("Feature 4 — Prepaid credits and admission control", 0, True),
        ("Generation claims its slot and debits a credit in a single transaction of conditional "
         "writes, so concurrent requests cannot double-spend. A semaphore caps concurrent "
         "generations at three; beyond that they queue and the player is shown their position.", 0, False),
    ], size=10.5)

    # ---- 8 implementation screenshots ------------------------------------
    clear_all(S[7])
    grid8 = [
        ("01-landing.png", "Landing page"),
        ("02-create-form.png", "Story creation form"),
        ("03-loading.png", "Generation in progress"),
        ("04-credits.png", "Credits and packs"),
    ]
    cw, ch = (W - Inches(0.9)) / 2, (H - Inches(1.5)) / 2
    for i, (fn, cap) in enumerate(grid8):
        cx = Inches(0.3) + (i % 2) * (cw + Inches(0.3))
        cy = Inches(0.35) + (i // 2) * (ch + Inches(0.5))
        rect = place(S[7], shot(fn), cx, cy, cw, ch - Inches(0.30), cap)
        if not Path(str(shot(fn))).exists():
            missing.append(fn)
        caption(S[7], cap, cx, rect[1] + rect[3] + Inches(0.04), cw)

    # ---- 9 demo screenshots ----------------------------------------------
    clear_all(S[8])
    caption(S[8], "Demo", Inches(0.3), Inches(0.10), W - Inches(0.6), size=17)
    grid9 = [
        ("05-reader.png", "Playing a generated story"),
        ("06-choices.png", "Branching choices"),
        ("07-explore.png", "Explore — published stories"),
        ("08-library.png", "Story page — ratings and comments"),
    ]
    gw = (W - Inches(0.9)) / 2
    gh = (H - Inches(1.35)) / 2
    for i, (fn, cap) in enumerate(grid9):
        cx = Inches(0.3) + (i % 2) * (gw + Inches(0.3))
        cy = Inches(0.72) + (i // 2) * (gh + Inches(0.18))
        rect = place(S[8], shot(fn), cx, cy, gw, gh - Inches(0.30), cap)
        if not Path(str(shot(fn))).exists():
            missing.append(fn)
        caption(S[8], cap, cx, rect[1] + rect[3] + Inches(0.04), gw, size=11)

    # ---- 10 flow diagram --------------------------------------------------
    clear_all(S[9])
    caption(S[9], "Generation Pipeline", Inches(0.3), Inches(0.15), W - Inches(0.6), size=17)
    place(S[9], CHARTS / "pipeline.png", Inches(0.3), Inches(0.85), W - Inches(0.6), H - Inches(1.4))

    # ---- 11 results -------------------------------------------------------
    set_title(S[10], "Results & Analysis")
    clear_pictures(S[10])
    set_body(S[10], [
        ("Output", 0, True),
        ("Deployed and operating at storyplex.app · 40 REST endpoints · 22 database tables · "
         "~9,300 lines of Python and ~4,100 of JavaScript", 0, False),
        ("Performance & cost", 0, True),
        ("Images per story 53 → 26 (51% fewer) · Memory baseline 750 MB → 129 MB · "
         "Cost per story ₹206 → ₹87 (58% lower) · 50 concurrent requests served in 147 ms", 0, False),
        ("Verification", 0, True),
        ("93 automated checks across four suites, all passing. Seven defects found and fixed during "
         "validation, including a double-spend race, a refund clawback gap, and unbounded "
         "concurrency that was causing out-of-memory crashes.", 0, False),
    ], size=11)
    place(S[10], CHARTS / "tests.png", Inches(2.6), Inches(4.35), Inches(4.8), Inches(2.55))

    # ---- 12 before / after -------------------------------------------------
    clear_all(S[11])
    caption(S[11], "Optimisation — measured before and after", Inches(0.3), Inches(0.12),
            W - Inches(0.6), size=16)
    place(S[11], CHARTS / "optimisation.png", Inches(0.25), Inches(0.75), W - Inches(0.5), Inches(3.1))
    place(S[11], CHARTS / "concurrency.png", Inches(1.9), Inches(4.05), Inches(6.2), Inches(3.1))

    # ---- 13 challenges ------------------------------------------------------
    set_title(S[12], "Challenges & Limitations")
    set_body(S[12], [
        ("Technical challenges", 0, True),
        ("1) Out-of-memory crashes in production. The ML background remover cost ~475 MB at import "
         "alone. Tuning it was not enough; replacing it with a chroma-key approach removed the "
         "dependency entirely and cut the baseline to 129 MB.", 0, False),
        ("2) Unbounded generation concurrency. Pipelines were dispatched with no admission control, "
         "so a fourth simultaneous generation exhausted memory and killed every in-flight story. "
         "Fixed with a semaphore and a player-visible queue.", 0, False),
        ("3) Concurrency in the billing path. Check-then-set on session status allowed two "
         "concurrent requests to both start a generation. Fixed by making the status claim and the "
         "credit debit conditional writes in one transaction.", 0, False),
        ("4) Provider model retirement. Gemini 2.5 models were withdrawn mid-project and generation "
         "failed in production. Model ids are now environment-overridable, and a pre-flight check "
         "verifies them.", 0, False),
        ("Limitations", 0, True),
        ("1) Generation takes ~10 minutes and is capped at three concurrent stories on the current "
         "instance size.", 0, False),
        ("2) Story quality varies with the premise and is not automatically assessed.", 0, False),
        ("3) Sprites occasionally retain a faint edge artefact from the source image's own "
         "anti-aliasing.", 0, False),
    ], size=10)

    # ---- 14 conclusion -------------------------------------------------------
    set_title(S[13], "Conclusion & Future Work")
    set_body(S[13], [
        ("Conclusion:", 0, True),
        ("StoryPlex demonstrates that a complete, playable visual novel — coherent branching story, "
         "consistent illustrated cast, and voiced dialogue — can be generated end to end from a "
         "short user premise. It is deployed and running as a real multi-user service with "
         "authentication, payments and social features. Beyond the generation pipeline itself, the "
         "project's substantive results are engineering ones: cost per story reduced 58% and memory "
         "footprint 83%, both by changing what the system generates rather than by degrading what "
         "the player receives.", 0, False),
        ("Future Work:", 0, True),
        ("1) Move generation to serverless jobs (Cloud Run) so concurrency scales beyond three and "
         "a deploy cannot interrupt an in-flight story.", 0, False),
        ("2) Persist generation progress to the database, enabling recovery of interrupted runs.", 0, False),
        ("3) Character reference-image conditioning for stronger visual consistency across scenes.", 0, False),
        ("4) Player-authored branches and community remixing of published stories.", 0, False),
        ("5) Multi-language generation and localisation of the reader interface.", 0, False),
    ], size=11)

    prs.save(str(OUT))
    return OUT, missing


if __name__ == "__main__":
    out, missing = build()
    print(f"  built: {out.relative_to(DOCS.parent)}  ({out.stat().st_size//1024} KB)")
    if missing:
        print(f"  placeholders for {len(missing)} screenshot(s) not yet supplied:")
        for m in missing:
            print(f"    · screenshots/{m}")
    else:
        print("  all screenshots present")
